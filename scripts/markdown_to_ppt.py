
"""
Convert markdown files to PowerPoint slides.
Extracts article title, focal points, date, and speaker notes to create presentation slides.
Usage:
    python markdown_to_ppt.py                           # Process all markdown files in updates/YYYY/MM/
    python markdown_to_ppt.py --from 2025-09-01         # From specific date onwards
    python markdown_to_ppt.py --pptx presentation.pptx   # Append to existing PPTX
    python markdown_to_ppt.py --pptx presentation.pptx --from 2025-09-01  # Combine options
"""
import os
import re
import argparse
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from urllib.request import urlopen, Request
from urllib.error import URLError
import json
import time

def fetch_og_image_from_url(url, max_retries=3, retry_delay=1):
    """
    Fetch the og:image meta tag from a given URL with retry policy.
    
    Args:
        url: URL to fetch og:image from
        max_retries: Maximum number of retry attempts (default: 3)
        retry_delay: Initial delay in seconds between retry attempts (default: 1)
                     Delays increment exponentially: 1s, 2s, 4s, etc.
    
    Returns:
        og:image URL if found, None otherwise
    """
    for attempt in range(max_retries):
        try:
            # Add a timeout and user-agent to avoid being blocked
            req = Request(
                url,
                headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
            )
            
            with urlopen(req, timeout=5) as response:
                html_content = response.read().decode('utf-8', errors='ignore')
                
                # Search for og:image meta tag
                og_image_match = re.search(
                    r'<meta\s+property=["\']og:image["\']\s+content=["\']([^"\']+)["\']|'
                    r'<meta\s+content=["\']([^"\']+)["\']\s+property=["\']og:image["\']',
                    html_content,
                    re.IGNORECASE
                )
                
                if og_image_match:
                    # Return the matched group (either group 1 or 2 depending on attribute order)
                    return og_image_match.group(1) or og_image_match.group(2)
                return None
        
        except (URLError, Exception) as e:
            # Retry on failure if attempts remain
            if attempt < max_retries - 1:
                # Exponential backoff: 1s, 2s, 4s, etc.
                wait_time = retry_delay * (2 ** attempt)
                time.sleep(wait_time)
                continue
            # Silently fail after all retries exhausted
            pass
    
    return None

def parse_markdown_file(file_path):
    """
    Parse markdown file and extract:
    - Article title (from # Article Title line at the top)
    - Article date (from ## Article Date section)
    - Article link (from ## Article Url section)
    - Focal points (from ## Article Content section with bullet points)
    - Speaker notes (from ### Speaker Notes section with Italian and English)
    - Article image (from ## Article Image section or **Image:** field, in multiple formats)
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Extract title from first "# " heading (main title)
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    title = title_match.group(1) if title_match else "Untitled"
    
    # Extract date from "## Article Date" section (next line after header)
    date_match = re.search(r'## Article Date\n(.+?)$', content, re.MULTILINE)
    article_date = date_match.group(1).strip() if date_match else ""
    
    # Extract link from "## Article Url" section (next line after header)
    link_match = re.search(r'## Article Url\n(.+?)$', content, re.MULTILINE)
    article_link = link_match.group(1).strip() if link_match else ""
    
    # Extract image - support multiple formats:
    # 1. ## Article Image followed by plain URL
    # 2. **Image:** [text](url) - markdown link
    # 3. **Image:** ![alt](url) - markdown image
    # 4. **Date:** YYYY-MM-DD at the top
    article_image = ""
    
    # Try format 1: ## Article Image section
    image_match = re.search(r'## Article Image\n(.+?)$', content, re.MULTILINE)
    if image_match:
        article_image = image_match.group(1).strip()
    
    # Try format 2 & 3: **Image:** at the top of file (before ## Article Content)
    if not article_image:
        image_match = re.search(r'^\*\*Image:\*\*\s*(.+?)$', content, re.MULTILINE)
        if image_match:
            article_image = image_match.group(1).strip()
    
    # Extract URL from various markdown formats
    if article_image:
        # Handle markdown link format: [text](url)
        if '[' in article_image and '](http' in article_image:
            url_match = re.search(r'\]\(([^)]+)\)', article_image)
            if url_match:
                article_image = url_match.group(1)
        # Handle markdown image format: ![alt](url)
        elif '![' in article_image and '](http' in article_image:
            url_match = re.search(r'\]\(([^)]+)\)', article_image)
            if url_match:
                article_image = url_match.group(1)
        # Handle plain URL (no additional extraction needed)
        # Extract just the URL if there's extra text
        url_match = re.search(r'(https?://[^\s\]]+)', article_image)
        if url_match:
            article_image = url_match.group(1)
    
    # Always fetch from og:image meta tag instead of using cached URL (with retry policy)
    if article_link:
        fetched_image = fetch_og_image_from_url(article_link, max_retries=3, retry_delay=1)
        if fetched_image:
            article_image = fetched_image
    
    # Extract focal points from "## Article Content" section
    # Support both bullet point format (• or -) and bold text format (**text**)
    focal_section = re.search(
        r'## Article Content\n((?:.|\n)+?)(?:^###|\Z)',
        content,
        re.MULTILINE
    )
    focal_points = []
    if focal_section:
        points_text = focal_section.group(1)
        for line in points_text.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
            # Support bullet point format (• or -)
            if line.startswith('• ') or line.startswith('- '):
                # Remove the bullet point prefix (either "• " or "- ")
                point = line[2:].strip()
                focal_points.append(point)
            # Support bold text format (**text**: description)
            elif line.startswith('**') and '**' in line[2:]:
                # Extract the entire line as a focal point
                # This preserves the format: **Title**: description
                focal_points.append(line)
    
    # Extract speaker notes from "### Speaker Notes" sections (both Italian and English)
    # Look for both "### **Speaker Notes (Italian)**:" and "### **Speaker Notes (English)**:"
    speaker_notes = ""

    # Find the positions of the speaker notes sections
    italian_header = "### **Speaker Notes (Italian)**:"
    english_header = "### **Speaker Notes (English)**:"

    italian_pos = content.find(italian_header)
    english_pos = content.find(english_header)

    # Extract Italian speaker notes
    if italian_pos != -1:
        italian_start = italian_pos + len(italian_header)
        # Find the end of Italian section (either English header or next main section)
        italian_end = english_pos if english_pos != -1 else content.find("\n## ", italian_start)
        if italian_end == -1:
            italian_end = len(content)
        italian_notes = content[italian_start:italian_end].strip()
        if italian_notes:
            speaker_notes += "**Italian Notes:**\n" + italian_notes + "\n\n"

    # Extract English speaker notes
    if english_pos != -1:
        english_start = english_pos + len(english_header)
        # Find the end of English section (next main section)
        english_end = content.find("\n## ", english_start)
        if english_end == -1:
            english_end = len(content)
        english_notes = content[english_start:english_end].strip()
        if english_notes:
            speaker_notes += "**English Notes:**\n" + english_notes

    speaker_notes = speaker_notes.strip()
    
    return {
        'title': title,
        'date': article_date,
        'link': article_link,
        'image': article_image,
        'focal_points': focal_points,
        'speaker_notes': speaker_notes
    }

def apply_segoe_ui_font(text_frame, font_size, bold=False, color=None, alignment=None):
    """Apply Segoe UI font styling to text frame.
    
    Args:
        text_frame: Text frame to style
        font_size: Font size in points
        bold: Whether to apply bold
        color: RGB color tuple
        alignment: Text alignment (PP_ALIGN.LEFT, PP_ALIGN.CENTER, etc.)
    """
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Segoe UI'
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
        if alignment:
            paragraph.alignment = alignment

def convert_markdown_to_pptx_format(text):
    """Convert markdown formatting to PowerPoint formatting.
    
    Converts **bold text** to bold text with proper PowerPoint formatting.
    
    Args:
        text: Text with markdown formatting
    
    Returns:
        Tuple of (plain_text, bold_positions) where bold_positions marks which parts should be bold
    """
    # Replace markdown bold markers with plain text
    return text.replace('**', '')

def parse_markdown_formatting(text_frame, text_with_markdown):
    """Parse markdown formatting in text and apply to PowerPoint text frame.
    
    Supports:
    - **bold text** -> bold
    - _italic text_ or *italic text* -> italic
    - `code text` -> monospace font (Courier New)
    - ~~strikethrough~~ -> strikethrough
    - [link text](url) -> hyperlink
    - Preserves line breaks and paragraph spacing
    
    Args:
        text_frame: PowerPoint text frame to populate
        text_with_markdown: Text with markdown formatting (may contain newlines)
    """
    # Clear existing text
    text_frame.clear()
    
    # Split text by newlines to preserve line breaks
    lines = text_with_markdown.split('\n')
    
    # Pattern to match markdown formatting: **bold**, _italic_, `code`, ~~strike~~, [link](url)
    pattern = r'(\*\*[^*]+\*\*|_[^_]+_|\*[^*]+\*|`[^`]+`|~~[^~]+~~|\[[^\]]+\]\([^)]+\)|[^*_`~\[](?:[^*_`~\[]*[^*_`~\[\s])?)'
    
    first_paragraph = True
    
    for line in lines:
        # Skip completely empty lines but preserve the spacing
        if not line.strip():
            if not first_paragraph:
                p = text_frame.add_paragraph()
                p.text = ""
            continue
        
        # Create paragraph for this line
        if first_paragraph:
            p = text_frame.paragraphs[0]
            first_paragraph = False
        else:
            p = text_frame.add_paragraph()
        
        # Find all markdown tokens in this line
        tokens = []
        last_end = 0
        
        for match in re.finditer(pattern, line):
            # Add any plain text before this match
            if match.start() > last_end:
                tokens.append(('plain', line[last_end:match.start()]))
            tokens.append(('markdown', match.group(0)))
            last_end = match.end()
        
        # Add any remaining plain text
        if last_end < len(line):
            tokens.append(('plain', line[last_end:]))
        
        # If no markdown found, just add the line as plain text
        if not tokens:
            tokens = [('plain', line)]
        
        # Process tokens and add runs to paragraph
        first_run = True
        for token_type, token in tokens:
            if not token.strip() and token_type == 'plain':
                continue
            
            if token_type == 'markdown':
                # Determine formatting based on markdown markers
                is_bold = token.startswith('**') and token.endswith('**')
                is_italic = (token.startswith('_') and token.endswith('_')) or (token.startswith('*') and token.endswith('*') and not token.startswith('**'))
                is_code = token.startswith('`') and token.endswith('`')
                is_strikethrough = token.startswith('~~') and token.endswith('~~')
                is_link = token.startswith('[') and '](' in token
                
                # Extract clean text (remove markdown markers)
                if is_bold:
                    clean_text = token[2:-2]
                elif is_italic:
                    clean_text = token[1:-1]
                elif is_code:
                    clean_text = token[1:-1]
                elif is_strikethrough:
                    clean_text = token[2:-2]
                elif is_link:
                    # Extract link text and URL: [text](url)
                    link_match = re.match(r'\[([^\]]+)\]\(([^)]+)\)', token)
                    clean_text = link_match.group(1) if link_match else token
                    link_url = link_match.group(2) if link_match else ""
                else:
                    clean_text = token
                
                # Add run to paragraph
                if first_run and not p.runs:
                    run = p.add_run()
                    first_run = False
                else:
                    run = p.add_run()
                
                run.text = clean_text
                
                # Apply formatting
                if is_bold:
                    run.font.bold = True
                if is_italic:
                    run.font.italic = True
                if is_code:
                    run.font.name = 'Courier New'
                    run.font.size = Pt(10)
                if is_strikethrough:
                    run.font.strikethrough = True
                if is_link:
                    # Add hyperlink
                    run.hyperlink.address = link_url
                    run.font.color.rgb = RGBColor(0, 0, 255)
                    run.font.underline = True
            else:
                # Plain text token
                if first_run and not p.runs:
                    run = p.add_run()
                    first_run = False
                else:
                    run = p.add_run()
                run.text = token

def format_focal_points_text(text_frame, focal_points):
    """Format focal points with bullet styling and Segoe UI font with markdown support.
    
    Args:
        text_frame: Text frame for focal points
        focal_points: List of focal point strings with markdown formatting
    """
    GITHUB_DARK = RGBColor(36, 41, 46)
    
    text_frame.clear()
    
    for i, point in enumerate(focal_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Add bullet point
        p.text = f"• "
        p.font.name = 'Segoe UI'
        p.font.size = Pt(18)
        p.font.color.rgb = GITHUB_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.level = 0
        p.word_wrap = True
        
        # Add focal point text with markdown formatting
        run = p.add_run()
        run.text = point
        run.font.name = 'Segoe UI'
        run.font.size = Pt(18)
        run.font.color.rgb = GITHUB_DARK
        
        # Parse markdown formatting in focal point text
        pattern = r'(\*\*[^*]+\*\*|_[^_]+_|\*[^*]+\*|`[^`]+`|~~[^~]+~~)'
        parts = re.split(pattern, point)
        
        # Clear the run text we just added
        p.clear()
        p.text = f"• "
        p.font.name = 'Segoe UI'
        p.font.size = Pt(18)
        p.font.color.rgb = GITHUB_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        
        # Add formatted parts
        for j, part in enumerate(parts):
            if not part:
                continue
                
            is_bold = part.startswith('**') and part.endswith('**')
            is_italic = (part.startswith('_') and part.endswith('_')) or (part.startswith('*') and part.endswith('*') and not part.startswith('**'))
            is_code = part.startswith('`') and part.endswith('`')
            is_strikethrough = part.startswith('~~') and part.endswith('~~')
            
            # Extract clean text
            if is_bold:
                clean_text = part[2:-2]
            elif is_italic:
                clean_text = part[1:-1]
            elif is_code:
                clean_text = part[1:-1]
            elif is_strikethrough:
                clean_text = part[2:-2]
            else:
                clean_text = part
            
            # Add run
            r = p.add_run()
            r.text = clean_text
            r.font.name = 'Segoe UI'
            r.font.size = Pt(18)
            r.font.color.rgb = GITHUB_DARK
            
            # Apply formatting
            if is_bold:
                r.font.bold = True
            if is_italic:
                r.font.italic = True
            if is_code:
                r.font.name = 'Courier New'
                r.font.size = Pt(16)
            if is_strikethrough:
                r.font.strikethrough = True

def create_cover_slide(title, image_url, prs, speaker_notes="", max_retries=3, retry_delay=1):
    """Create a cover slide with article title and image as full background with retry policy.
    
    Args:
        title: Article title
        image_url: URL to the article image
        prs: Presentation object
        speaker_notes: Speaker notes text to add to the slide
        max_retries: Maximum number of retry attempts for image download (default: 3)
        retry_delay: Delay in seconds between retry attempts (default: 1)
    
    Returns:
        Presentation object with cover slide added
    """
    from io import BytesIO
    
    # Add blank slide
    layout = prs.slide_layouts[3]  # Blank layout
    slide = prs.slides.add_slide(layout)
    
    # Add black background first
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add image as full background if URL is provided (with retry policy)
    image_added = False
    if image_url and image_url != "https://github.blog/changelog/":
        for attempt in range(max_retries):
            try:
                # Download image from URL
                req = Request(
                    image_url,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                )
                response = urlopen(req, timeout=5)
                image_data = BytesIO(response.read())
                
                # Add image as full background (fill entire slide)
                slide.shapes.add_picture(image_data, 0, 0, width=prs.slide_width, height=prs.slide_height)
                image_added = True
                break
            except Exception as e:
                # Retry on failure if attempts remain
                if attempt < max_retries - 1:
                    # Exponential backoff: 1s, 2s, 4s, etc.
                    wait_time = retry_delay * (2 ** attempt)
                    time.sleep(wait_time)
                    continue
                # If image download fails after all retries, continue with black background
                pass
    
    # Add semi-transparent dark overlay to ensure text readability over the image
    if image_added:
        overlay = slide.shapes.add_shape(1, 0, Inches(0), prs.slide_width, Inches(1.3))  # 1 = rectangle
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.4  # 40% opacity
        overlay.line.color.rgb = RGBColor(0, 0, 0)
        overlay.line.width = 0
    
    # Add title as text box overlay with white color (full slide width)
    title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), prs.slide_width - Inches(1), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.name = 'Segoe UI'
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)  # White color
    
    # Add speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.clear()
        parse_markdown_formatting(notes_text_frame, speaker_notes)
    
    return prs

def create_single_slide(markdown_file, prs=None):
    """Create a single comprehensive slide per markdown file.
    
    Slide includes:
    - Article title (32pt Segoe UI Semibold, black)
    - Focal points list (28pt Segoe UI, black, with bullets)
    - Date footer in bottom right corner (18pt Segoe UI, light gray)
    - Speaker notes with date, link, and multi-language notes
    
    Uses layout placeholders (0 for title, 1 for content) if available,
    otherwise creates text boxes dynamically.
    
    Args:
        markdown_file: Path to markdown file
        prs: Optional existing Presentation object
    
    Returns:
        Presentation object with new slide added
    """
    
    # Parse markdown
    data = parse_markdown_file(markdown_file)
    
    # Load existing presentation or create new one
    if prs is None:
        prs = Presentation()
        # Use widescreen format (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    speaker_notes_text = f"Date: {data['date']} - Link: {data['link']}\n\n{data['speaker_notes']}"
    
    # Create cover slide with title and image (with retry policy for image download)
    prs = create_cover_slide(data['title'], data['image'], prs, speaker_notes_text, max_retries=3, retry_delay=1)
    
    # GitHub colors
    GITHUB_DARK = RGBColor(36, 41, 46)
    GITHUB_LIGHT_GRAY = RGBColor(128, 128, 128)
    
    # Try to use layout with placeholders (index 3 typically has title and content)
    # Otherwise use blank layout (index 6)
    layout_index = 3 if len(prs.slide_layouts) > 3 else 6
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    
    # Remove all placeholders except the first one
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx != 0:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Use text box for title that fills slide width and starts from left
    # Title takes up full width (from left edge to right edge) with no margins
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0.3), prs.slide_width, Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = data['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = GITHUB_DARK
    title_p.alignment = PP_ALIGN.LEFT  # Align to left
    
    # Set content using text box (no content placeholder)
    # Calculate width with margins (0.5 left + 0.5 right = 1 inch total margin)
    content_width = prs.slide_width - Inches(1)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), content_width, Inches(2.3))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    format_focal_points_text(content_frame, data['focal_points'])
    
    # Add date footer in bottom right corner (18pt Segoe UI, light gray)
    date_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.9), Inches(2.2), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.word_wrap = True
    date_p = date_frame.paragraphs[0]
    date_p.text = data['date']
    date_p.font.name = 'Segoe UI'
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = GITHUB_LIGHT_GRAY
    date_p.alignment = PP_ALIGN.RIGHT
    
    # Add speaker notes with date, link, and multi-language notes
    #speaker_notes_text = f"Date: {data['date']} - Link: {data['link']}\n\n{data['speaker_notes']}"
    
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.clear()
    
    # Parse markdown formatting in speaker notes
    parse_markdown_formatting(notes_text_frame, speaker_notes_text)
    
    return prs

def filter_files_by_date(files, date_from=None, date_to=None):
    """Filter markdown files by date range."""
    filtered = []
    
    for file_path in files:
        # Extract date from filename (format: YYYY-MM-DD-*.md)
        match = re.match(r'(\d{4}-\d{2}-\d{2})', file_path.name)
        if not match:
            continue
        
        file_date = match.group(1)
        
        try:
            file_datetime = datetime.strptime(file_date, '%Y-%m-%d')
            
            # Check if file is within date range
            if date_from:
                from_datetime = datetime.strptime(date_from, '%Y-%m-%d')
                if file_datetime < from_datetime:
                    continue
            
            if date_to:
                to_datetime = datetime.strptime(date_to, '%Y-%m-%d')
                if file_datetime > to_datetime:
                    continue
            
            filtered.append(file_path)
        except ValueError:
            continue
    
    return filtered

def main():
    """Process markdown files with optional date filtering."""
    
    parser = argparse.ArgumentParser(
        description='Convert markdown files to PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python markdown_to_ppt.py                                      # Create updates.pptx with all files
  python markdown_to_ppt.py --from 2025-09-01                   # From date onwards
  python markdown_to_ppt.py --from 2025-09-01 --to 2025-09-30   # Date range
  python markdown_to_ppt.py --output custom.pptx                # Specify output filename
  python markdown_to_ppt.py --append updates.pptx               # Append to existing PPTX
        """
    )
    
    parser.add_argument(
        '--from',
        dest='date_from',
        type=str,
        help='Start date (YYYY-MM-DD) to filter markdown files'
    )
    parser.add_argument(
        '--to',
        dest='date_to',
        type=str,
        help='End date (YYYY-MM-DD) to filter markdown files'
    )
    parser.add_argument(
        '--output',
        type=str,
        default='updates.pptx',
        help='Output PPTX filename (default: updates.pptx)'
    )
    parser.add_argument(
        '--append',
        type=str,
        help='Path to existing PPTX file to append slides to'
    )
    
    args = parser.parse_args()
    
    current_dir = Path(os.path.dirname(os.path.abspath(__file__)))
    updates_dir = current_dir / 'updates'
    
    if not updates_dir.exists():
        print(f"✗ Updates directory not found: {updates_dir}")
        return
    
    # Recursively search for markdown files in YYYY/MM/ subdirectories
    markdown_files = []
    for year_dir in sorted(updates_dir.glob('[0-9][0-9][0-9][0-9]')):
        if year_dir.is_dir():
            for month_dir in sorted(year_dir.glob('[0-9][0-9]')):
                if month_dir.is_dir():
                    markdown_files.extend(month_dir.glob('*.md'))
    
    # Sort all files by filename (which includes date)
    markdown_files = sorted(markdown_files)
    
    if not markdown_files:
        print("✗ No markdown files found in updates directory")
        return
    
    # Filter by date if specified
    if args.date_from or args.date_to:
        markdown_files = filter_files_by_date(markdown_files, args.date_from, args.date_to)
        
        if args.date_from and args.date_to:
            print(f"📅 Filtering files from {args.date_from} to {args.date_to}")
        elif args.date_from:
            print(f"📅 Filtering files from {args.date_from} onwards")
        elif args.date_to:
            print(f"📅 Filtering files up to {args.date_to}")
    
    if not markdown_files:
        print("✗ No markdown files found in the specified date range")
        return
    
    # Create pptx output folder if it doesn't exist
    pptx_dir = current_dir / 'pptx'
    pptx_dir.mkdir(exist_ok=True)
    
    # Load existing PPTX if provided with --append
    shared_prs = None
    if args.append:
        append_path = Path(args.append)
        if append_path.exists():
            try:
                shared_prs = Presentation(str(append_path))
                print(f"✓ Loaded existing PPTX: {args.append}")
            except Exception as e:
                print(f"✗ Error loading PPTX {args.append}: {e}")
                return
        else:
            print(f"✗ File not found: {args.append}")
            return
    else:
        # Create new presentation for consolidated output with widescreen format
        shared_prs = Presentation()
        shared_prs.slide_width = Inches(13.333)
        shared_prs.slide_height = Inches(7.5)
    
    # Generate unique filename based on date filters if using default output name
    if args.output == 'updates.pptx' and (args.date_from or args.date_to):
        # Build filename from date range
        if args.date_from and args.date_to:
            output_filename = f"updates_{args.date_from}_to_{args.date_to}.pptx"
        elif args.date_from:
            output_filename = f"updates_from_{args.date_from}.pptx"
        else:  # args.date_to only
            output_filename = f"updates_to_{args.date_to}.pptx"
    else:
        output_filename = args.output
    
    # Construct full output path in pptx folder
    output_pptx = str(pptx_dir / output_filename)
    
    processed_count = 0
    for md_file in markdown_files:
        try:
            # Add all slides to the consolidated PPTX
            shared_prs = create_single_slide(str(md_file), prs=shared_prs)
            print(f"✓ Added slide from: {md_file.name}")
            processed_count += 1
        except Exception as e:
            print(f"✗ Error processing {md_file.name}: {e}")
    
    # Save consolidated PPTX
    if processed_count > 0:
        shared_prs.save(output_pptx)
        print(f"\n✓ Saved consolidated PPTX: {output_pptx}")
    else:
        print(f"\n✗ No slides were created")
    
    print(f"✓ Successfully processed {processed_count} markdown file(s)")

if __name__ == '__main__':
    main()
